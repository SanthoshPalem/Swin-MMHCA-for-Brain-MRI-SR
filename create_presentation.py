
import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# --- Configuration ---
PPT_FILE_NAME = "Swin_MMHCA_Presentation.pptx"
RESULTS_DIR = "results"
IMG_DIR = "presentation_assets"
SRM_LOGO_URL = "https://srmap.edu.in/wp-content/themes/srm-theme/assets/images/srm-logo.png"
BASELINE_ARCH_URL = "https://raw.githubusercontent.com/lilygeorgescu/MHCA/main/imgs/overview.png"
QUALITATIVE_IMG_NAME = "comparison_sample_1.png" # Assumes this exists in 'results'
SRM_LOGO_PATH = os.path.join(IMG_DIR, "srmap_logo.png")
BASELINE_ARCH_PATH = os.path.join(IMG_DIR, "baseline_arch.png")
QUALITATIVE_IMG_PATH = os.path.join(RESULTS_DIR, QUALITATIVE_IMG_NAME)

# --- Helper Functions ---
def download_image(url, path):
    """Downloads an image from a URL and saves it."""
    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()
        with open(path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        print(f"Downloaded {os.path.basename(path)}")
        return True
    except requests.exceptions.RequestException as e:
        print(f"Error downloading {url}: {e}")
        return False

def add_title_slide(prs):
    """Adds the main title slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Swin-MMHCA: Integrating Swin Transformers with Multi-Modal Attention for Medical Image Super-Resolution"
    subtitle.text = "Authors: Ospari Jagadeesh, Palem Santhosh, Padyala Chakravathi\n" \
                    "Guide: Dr. Syed Sameen Ahmad Rizvi\n" \
                    "Department of Computer Science & Engineering, SRM University-AP\n" \
                    "December 2025"
    
    # Add Logo
    if os.path.exists(SRM_LOGO_PATH):
        slide.shapes.add_picture(SRM_LOGO_PATH, Inches(8.5), Inches(0.2), height=Inches(1.0))

def add_bullet_slide(prs, title_text, bullet_points):
    """Adds a standard title and content slide with bullet points."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = title_text
    tf = body.text_frame
    tf.clear()
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        if i == 0:
            p.font.bold = True
        if point.startswith('\t'):
            p.text = point.lstrip('\t')
            p.level = 1
            p.font.bold = False

def add_lit_review_slide(prs):
    title = "Literature Review & Novelty"
    points = [
        "Foundation 1: EDSR + MMHCA (Baseline)",
        "\tBased on Georgescu et al. (WACV 2023).",
        "\tUses a CNN backbone (EDSR) and the MMHCA module for data fusion.",
        "Foundation 2: Swin Transformer",
        "\tA Vision Transformer that uses shifted window self-attention.",
        "\tEfficiently captures both local and global features.",
        "Our Novelty",
        "\tWe replace the CNN backbone with a Swin Transformer.",
        "\tThis creates the Swin-MMHCA model to better capture global context."
    ]
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for point in points:
        p = tf.add_paragraph()
        p.text = point.lstrip('\t')
        if not point.startswith('\t'):
            p.level = 0
            p.font.bold = True
        else:
            p.level = 1
            p.font.bold = False

    if os.path.exists(BASELINE_ARCH_PATH):
        # Add baseline architecture image to the right
        slide.shapes.add_picture(BASELINE_ARCH_PATH, Inches(5.5), Inches(4.0), width=Inches(4.3))


def add_architecture_slide(prs):
    """Adds the slide for the proposed architecture with a placeholder for the Mermaid diagram."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Proposed Architecture: Swin-MMHCA"
    
    tf = body.text_frame
    tf.clear()
    p1 = tf.add_paragraph()
    p1.text = "This diagram shows the flow of data through our model:"
    p1.level = 0
    
    # Placeholder for Mermaid diagram
    left, top, width, height = Inches(1.5), Inches(2.5), Inches(7), Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf2 = textbox.text_frame
    tf2.word_wrap = True
    p2 = tf2.add_paragraph()
    p2.text = "NOTE: Please replace this text box with an image of the Mermaid diagram below."
    p2.font.color.rgb = RGBColor(255, 0, 0)
    p2.font.bold = True

    p3 = tf2.add_paragraph()
    p3.text = (
        'graph TD;\n'
        '    subgraph Inputs\n'
        '        A1[LR Input 1 (T1)] --> B1(CNN Encoder);\n'
        '        A2[LR Input 2 (T2)] --> B2(CNN Encoder);\n'
        '        A3[LR Input 3 (PD)] --> B3(CNN Encoder);\n'
        '    end\n'
        '    B1 & B2 & B3 --> C(Concatenate & Add Positional Encoding);\n'
        '    C --> E{Swin Transformer <br> Deep Feature Extraction};\n'
        '    E --> F(Reshape);\n'
        '    F --> G(MMHCA Module <br> Attention Fusion);\n'
        '    G --> H(CNN Upsampling Decoder);\n'
        '    H --> I[HR Output Image];'
    )
    p3.font.size = Pt(10)

def add_results_table_slide(prs):
    """Adds the slide with the quantitative results table."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Quantitative Results (4x Super-Resolution)"
    
    # Add bullet points for metrics
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Metrics Used:"
    p1 = tf.add_paragraph()
    p1.text = "PSNR (Peak Signal-to-Noise Ratio): Higher is better."
    p1.level = 1
    p2 = tf.add_paragraph()
    p2.text = "SSIM (Structural Similarity Index): Higher is better."
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "LPIPS (Learned Perceptual Image Patch Similarity): Lower is better."
    p3.level = 1
    
    # Add table
    table_data = [
        ['Model', 'PSNR (dB) \u2191', 'SSIM \u2191', 'LPIPS \u2193'],
        ['EDSR_Nav (Baseline)', '23.67', '0.135', '0.813'],
        ['SwinMMHCA (Ours)', '33.90', '0.768', '0.236']
    ]
    rows, cols = 3, 4
    left, top, width, height = Inches(1.5), Inches(4.0), Inches(7.0), Inches(1.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = table_data[r][c]
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            if r == 0:
                para.font.bold = True
            if r == 2:
                para.font.bold = True

def add_qualitative_slide(prs):
    """Adds the slide with the visual comparison image."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Qualitative Results (Visual Comparison)"
    
    # Add image
    if os.path.exists(QUALITATIVE_IMG_PATH):
        slide.shapes.add_picture(QUALITATIVE_IMG_PATH, Inches(0.5), Inches(2.0), width=Inches(9.0))
        # Add caption
        left, top, width, height = Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = "Our model produces visibly sharper images with better-defined structural details."
    else:
        body = slide.placeholders[1]
        body.text = f"Image not found: {QUALITATIVE_IMG_PATH}\nRun visualization to generate it."

def main():
    # --- Setup ---
    if not os.path.exists(IMG_DIR):
        os.makedirs(IMG_DIR)

    print("Downloading required image assets...")
    download_image(SRM_LOGO_URL, SRM_LOGO_PATH)
    download_image(BASELINE_ARCH_URL, BASELINE_ARCH_PATH)
    
    # --- Presentation Creation ---
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating Slide 1: Title")
    add_title_slide(prs)

    print("Creating Slide 2: Introduction & Motivation")
    add_bullet_slide(prs, "Introduction & Motivation", [
        "What is Super-Resolution?",
        "\tReconstructing High-Resolution (HR) images from Low-Resolution (LR) inputs.",
        "\tMatters in Medicine: Enhances diagnostic clarity without costly hardware or long scan times.",
        "The Problem with Existing Methods:",
        "\tTraditional CNNs struggle with long-range dependencies due to their local receptive fields.",
        "Our Objective:",
        "\tDesign a novel architecture, Swin-MMHCA, using a Swin Transformer to better capture global context."
    ])
    
    print("Creating Slide 3: Literature Review & Novelty")
    add_lit_review_slide(prs)

    print("Creating Slide 4: Dataset and Preprocessing")
    add_bullet_slide(prs, "Dataset and Preprocessing", [
        "Dataset: IXI (Information eXtraction from Images)",
        "\t~600 MRI scans from healthy subjects.",
        "\tModalities Used: T1-weighted, T2-weighted, and Proton Density (PD).",
        "Preprocessing Pipeline:",
        "\tCentral 2D slice extracted from 3D MRI volume.",
        "\tLR input created via bicubic downsampling.",
        "Data Split:",
        "\tTraining: 500 subjects, Test: 70 subjects"
    ])

    print("Creating Slide 5: Proposed Architecture")
    add_architecture_slide(prs)

    print("Creating Slide 6: How It Works")
    add_bullet_slide(prs, "How It Works: Architecture Breakdown", [
        "1. CNN Encoders:",
        "\tShallow conv layers create initial feature maps from each modality.",
        "2. Positional Encoding & Concatenation:",
        "\tFeatures are combined and spatial position info is added.",
        "3. Swin Transformer:",
        "\tThe core of the model. A powerful deep feature extractor for global context.",
        "4. MMHCA Module:",
        "\tAttention mechanism that fuses the features from the transformer.",
        "5. Upsampling Decoder:",
        "\tTransposed convolution layers that upscale features to the final HR image."
    ])

    print("Creating Slide 7: Training Details")
    add_bullet_slide(prs, "Training Details", [
        "Framework: PyTorch",
        "Optimizer: Adam",
        "Loss Function: L1 Loss (Mean Absolute Error)",
        "Epochs: 50",
        "Batch Size: 4",
        "Hardware: NVIDIA GeForce RTX 3050 Laptop GPU",
        "Total Training Time: ~3 hours, 5 minutes"
    ])
    
    print("Creating Slide 8: Quantitative Results")
    add_results_table_slide(prs)
    
    print("Creating Slide 9: Qualitative Results")
    add_qualitative_slide(prs)

    print("Creating Slide 10: Conclusion & Future Work")
    add_bullet_slide(prs, "Conclusion & Future Work", [
        "Conclusion:",
        "\tSuccessfully designed and implemented Swin-MMHCA.",
        "\tAchieved significant performance improvement over the CNN-based baseline.",
        "\tValidated that transformers are highly effective for this task.",
        "Future Work:",
        "\tConduct extensive hyperparameter tuning.",
        "\tExplore more advanced decoder modules.",
        "\tCollaborate with medical professionals for clinical evaluation."
    ])

    # --- Save Presentation ---
    prs.save(PPT_FILE_NAME)
    print(f"\nSuccessfully generated presentation: {os.path.abspath(PPT_FILE_NAME)}")

if __name__ == "__main__":
    main()
